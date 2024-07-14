import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from googlesearch import search
import requests
from bs4 import BeautifulSoup
from io import BytesIO
import os
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

title = "Alpine Adventures: Experience Switzerland, Japan, and Argentina"
text = """
Based on this information:Sure, based on the information provided, here are three temperate countries to travel to in January, one from each of three different continents:

1. France (Europe)
Attraction 1: Eiffel Tower: The iconic symbol of Paris and one of the most recognizable structures in the world, the Eiffel Tower offers stunning views of the city from its various viewing platforms.
Attraction 2: Louvre Museum: Home to thousands of works of art, including the Mona Lisa and the Venus de Milo, the Louvre is one of the largest and most famous museums in the world.
Attraction 3: Versailles Palace: A UNESCO World Heritage site, the Palace of Versailles is a magnificent example of French Baroque architecture and contains numerous ornate rooms, gardens, and fountains.
2. Argentina (South America)
Attraction 1: Iguazu Falls: Located on the border of Argentina and Brazil, Iguazu Falls is one of the largest and most impressive waterfall systems in the world, featuring hundreds of cascades spread out over nearly two miles.
Attraction 2: Perito Moreno Glacier: A massive ice formation located in Los Glaciares National Park in southern Argentina, Perito Moreno Glacier is a popular destination for hiking, ice climbing, and boat tours.
Attraction 3: Buenos Aires: The capital city of Argentina, Buenos Aires is known for its vibrant culture, delicious cuisine, and historic architecture, including the famous Teatro Colón opera house and the colorful neighborhoods of La Boca and San Telmo.
3. Japan (Asia)
Attraction 1: Mount Fuji: The highest mountain in Japan and a sacred site for many Japanese people, Mount Fuji is a popular destination for hiking and climbing, especially during the warmer months.
Attraction 2: Tokyo: The bustling capital city of Japan, Tokyo is home to numerous attractions, including the famous Shibuya Crossing, the ancient Senso-ji Temple, and the modern Akihabara district, known for its electronics stores and anime shops.
Attraction 3: Kyoto: A former imperial capital of Japan, Kyoto is known for its traditional temples, gardens, and tea ceremonies, as well as its geisha district and the famous Fushimi Inari Shrine, with its thousands of vermilion torii gates.if I am creating a brochure based on countries to travel to, what are 3 formal and professional names I can give to the travel brochure to attract tourism
"""
def generate_slides(text,  theme='summer', title='My Trip'):
    ## extract country names
    countries = re.findall(r'^\d+\.\s+([A-Za-z\s]+)\s+\(.*?\)', text, re.MULTILINE)
    ## extract number of slides

    title = title.strip('"')
    number = len(countries) #number of slides
    print(number)
    matches = re.findall(r'Attraction \d+: (.*?)(?=\.)', text)
    nested_list = [matches[i:i+3] for i in range(0, len(matches), 3)]
    print(nested_list)


    base_url = "https://stock.adobe.com/search?k="



    urls = [base_url + country + "%20travel" for country in countries]

    def extract_image_url(url):
        try:
            # Send an HTTP GET request to the URL
            response = requests.get(url)
            response.raise_for_status()  # Raise an exception for invalid responses

            # Parse the HTML content of the webpage
            soup = BeautifulSoup(response.content, 'html.parser')

            # Find all image tags (img) with non-empty src attributes
            img_tags = soup.find_all('img', src=True)

            # Extract the URLs of the images (src attribute of the img tags)
            for img_tag in img_tags:
                src = img_tag['src']
                # Check if the URL contains '.jpg' or '.jpeg' (assuming images are JPEG format)
                if '.jpg' in src or '.jpeg' in src:
                    # Exclude URLs containing 'adobe-logo' (assumed to be stock Adobe logos)
                    if 'adobe-logo' not in src:
                        # Return the URL of the image
                        return src

        except Exception as e:
            print(f"Error occurred while processing {url}: {e}")

        # Return None if no valid image URL is found
        return None

    image_urls = []
    for url in urls:
        image_url = extract_image_url(url)
        if image_url:
            image_urls.append(image_url)

    # Print the list of image URLs
    for image_url in image_urls:
        print(image_url)

    num_slides = number
    theme = theme.lower()
    prs = Presentation()

    ## slides formatting for summer
    if theme == 'summer':
        title_slide_layout = prs.slide_layouts[0]  # Layout index 0 is the title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)

        # Add background image
        try:
            bg_image_path = "themes/summer/summer-intro.png"  # Path to your local image
            bg_image = title_slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            bg_image.zorder = -10  # Set the zorder to place it behind other shapes
        except Exception as e:
            print(f"Failed to load background image: {e}")

        left = Inches(1)  # Adjust position as needed
        top = Inches(1)  # Adjust position as needed
        width = Inches(8)  # Adjust size as needed
        height = Inches(3)  # Adjust size as needed
        txBox_intro = title_slide.shapes.add_textbox(left, top, width, height)
        tf_intro = txBox_intro.text_frame
        p_intro = tf_intro.add_paragraph()
        p_intro.text = title
        tf_intro.word_wrap = True
        p_intro.font.bold = True
        p_intro.font.size = Pt(48)
        p_intro.font.name = 'Lucida Handwriting'
        p_intro.font.color.rgb = RGBColor(27,73,94)
        for paragraph in tf_intro.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

        # Add slides with country names on each slide
        for i in range(num_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use the layout for title slides (layout index 5)

            try:
                bg_image_path = "themes/summer/summer-bg.png"  # Path to your local image
                bg_image = slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
                bg_image.zorder = -10  # Set the zorder to place it behind other shapes

            except Exception as e:
                print(f"Failed to load background image: {e}")

            # Add title as a text box with similar structure to attractions with pointers
            left = Inches(6)
            top = Inches(1)
            width = Inches(3)
            height = Inches(1)
            txBox_title = slide.shapes.add_textbox(left, top, width, height)
            tf_title = txBox_title.text_frame
            tf_title.word_wrap = True

            p_title = tf_title.add_paragraph()
            p_title.text = countries[i]
            p_title.font.name = 'Didact Gothic'
            p_title.font.size = Pt(40)
            p_title.font.color.rgb = RGBColor(27, 73, 94)

            

            # Set the zorder of the title shape to ensure it's above the background image
            txBox_title.zorder = 10

            # Add attractions with pointers
            left = Inches(0.8)
            top = Inches(1.25)
            width = Inches(4.5)
            height = Inches(5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True

            # Calculate the width of the text box for the country title
            title_text_width = Inches(7.62)

            # Calculate the length of the text box for the body text
            body_text_length = Inches(11.3)

            # Resize the body text font size if it exceeds the width
            for j, attraction in enumerate(nested_list[i]):
                p = tf.add_paragraph()
                p.text = f"{j+1}. {attraction}"
                p.font.size = Pt(16) if txBox.width > body_text_length else Pt(20)  # Adjust font size
                p.space_before = Pt(14)  # Adjust spacing between each attraction
                p.font.name = 'Didact Gothic'
                p.font.color.rgb = RGBColor(27, 73, 94)

            # Add image from URL
            print(i, image_urls)
            image = image_urls[i]
            try:
                response = requests.get(image)
                if response.status_code == 200:
                    image_stream = BytesIO(response.content)
                    slide.shapes.add_picture(image_stream, Inches(6), Inches(2.5), width=Inches(3), height=Inches(3))
            except Exception as e:
                print(f"Failed to load image from {image}: {e}")

        end_slide_layout = prs.slide_layouts[0]  # Layout index 0 is the title slide layout
        end_slide = prs.slides.add_slide(end_slide_layout)


        # Add background image
        try:
            bg_image_path = "themes/summer/summer-conclusion.png"  # Path to your local image
            bg_image = end_slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            bg_image.zorder = -10  # Set the zorder to place it behind other shapes
        except Exception as e:
            print(f"Failed to load background image: {e}")

        left = Inches(3)  # Adjust position as needed
        top = Inches(3)  # Adjust position as needed
        width = Inches(8)  # Adjust size as needed
        height = Inches(1)  # Adjust size as needed
        txBox_end = end_slide.shapes.add_textbox(left, top, width, height)
        tf_end = txBox_end.text_frame
        p_end = tf_end.add_paragraph()
        p_end.text = 'Thank You!'
        p_end.font.bold = True
        p_end.font.size = Pt(40)
        p_end.font.name = 'Lucida Handwriting'
        p_end.font.color.rgb = RGBColor(27, 73, 94)



    ## slides formatting for winter
    if theme == 'winter':
        title_slide_layout = prs.slide_layouts[0]  # Layout index 0 is the title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)

        # Add background image
        try:
            bg_image_path = "themes/winter/winter-intro.png"  # Path to your local image
            bg_image = title_slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            bg_image.zorder = -10  # Set the zorder to place it behind other shapes
        except Exception as e:
            print(f"Failed to load background image: {e}")

        left = Inches(1)  # Adjust position as needed
        top = Inches(1)  # Adjust position as needed
        width = Inches(8)  # Adjust size as needed
        height = Inches(3)  # Adjust size as needed
        txBox_intro = title_slide.shapes.add_textbox(left, top, width, height)
        tf_intro = txBox_intro.text_frame
        p_intro = tf_intro.add_paragraph()
        p_intro.text = title
        tf_intro.word_wrap = True
        p_intro.font.bold = True
        p_intro.font.size = Pt(48)
        p_intro.font.name = 'ADLaM Display'
        p_intro.font.color.rgb = RGBColor(37, 64, 97)
        for paragraph in tf_intro.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

        for i in range(num_slides):
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use the layout for title slides (layout index 5)

                try:
                    bg_image_path = "themes/winter/winter-bg.png"  # Path to your local image
                    bg_image = slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
                    bg_image.zorder = -10  # Set the zorder to place it behind other shapes

                except Exception as e:
                    print(f"Failed to load background image: {e}")


                # Add title as a text box with similar structure to attractions with pointers
                left = Inches(4)
                top = Inches(0.1)
                width = Inches(9)
                height = Inches(1)
                txBox_title = slide.shapes.add_textbox(left, top, width, height)
                tf_title = txBox_title.text_frame
                tf_title.word_wrap = True
                p_title = tf_title.add_paragraph()
                p_title.text = countries[i]
                p_title.font.name = 'ADLaM Display'
                p_title.font.size = Pt(40)
                p_title.font.color.rgb = RGBColor(255,255,255)

                

                # Set the zorder of the title shape to ensure it's above the background image
                txBox_title.zorder = 10

                # Add attractions with pointers
                left = Inches(0.8)
                top = Inches(1.25)
                width = Inches(4.5)
                height = Inches(5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True

                # Calculate the length of the text box for the body text
                body_text_length = Inches(11.3)

                # Resize the body text font size if it exceeds the width
                for j, attraction in enumerate(nested_list[i]):
                    p = tf.add_paragraph()
                    p.text = f"{j+1}. {attraction}"
                    p.font.size = Pt(16) if txBox.width > body_text_length else Pt(20)  # Adjust font size
                    p.space_before = Pt(14)  # Adjust spacing between each attraction
                    p.font.color.rgb = RGBColor(255,255,255)

                # Add image from URL
                image = image_urls[i]
                try:
                    response = requests.get(image)
                    if response.status_code == 200:
                        image_stream = BytesIO(response.content)
                        slide.shapes.add_picture(image_stream, Inches(6), Inches(2.5), width=Inches(3), height=Inches(3))
                except Exception as e:
                    print(f"Failed to load image from {image}: {e}")

        end_slide_layout = prs.slide_layouts[0]  # Layout index 0 is the title slide layout
        end_slide = prs.slides.add_slide(end_slide_layout)


        # Add background image
        try:
            bg_image_path = "themes/winter/winter-conclusion.png"  # Path to your local image
            bg_image = end_slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            bg_image.zorder = -10  # Set the zorder to place it behind other shapes
        except Exception as e:
            print(f"Failed to load background image: {e}")

        left = Inches(3)  # Adjust position as needed
        top = Inches(3)  # Adjust position as needed
        width = Inches(8)  # Adjust size as needed
        height = Inches(1)  # Adjust size as needed
        txBox_end = end_slide.shapes.add_textbox(left, top, width, height)
        tf_end = txBox_end.text_frame
        p_end = tf_end.add_paragraph()
        p_end.text = 'Thank You!'
        p_end.font.bold = True
        p_end.font.size = Pt(55)
        p_end.font.name = 'ADLaM Display'
        p_end.font.color.rgb = RGBColor(255, 255, 255)


    ## slides formatting for fall
    if theme == 'fall':
        title_slide_layout = prs.slide_layouts[0]  # Layout index 0 is the title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)

        # Add background image
        try:
            bg_image_path = "themes/fall/fall-intro.png"  # Path to your local image
            bg_image = title_slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            bg_image.zorder = -10  # Set the zorder to place it behind other shapes
        except Exception as e:
            print(f"Failed to load background image: {e}")

        left = Inches(1)  # Adjust position as needed
        top = Inches(1)  # Adjust position as needed
        width = Inches(8)  # Adjust size as needed
        height = Inches(3)  # Adjust size as needed
        txBox_intro = title_slide.shapes.add_textbox(left, top, width, height)
        tf_intro = txBox_intro.text_frame
        p_intro = tf_intro.add_paragraph()
        p_intro.text = title
        tf_intro.word_wrap = True
        p_intro.font.bold = True
        p_intro.font.size = Pt(48)
        p_intro.font.name = 'COPPERPLATE GOTHIC'
        p_intro.font.color.rgb = RGBColor(139,50,28)
        for paragraph in tf_intro.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
        
        for i in range(num_slides):
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use the layout for title slides (layout index 5)

                try:
                    bg_image_path = "themes/fall/fall-bg.png"  # Path to your local image
                    bg_image = slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
                    bg_image.zorder = -10  # Set the zorder to place it behind other shapes

                except Exception as e:
                    print(f"Failed to load background image: {e}")



                # Add title as a text box with similar structure to attractions with pointers
                left = Inches(4)
                top = Inches(0.1)
                width = Inches(9)
                height = Inches(1)
                txBox_title = slide.shapes.add_textbox(left, top, width, height)
                tf_title = txBox_title.text_frame
                tf_title.word_wrap = True
                p_title = tf_title.add_paragraph()
                p_title.text = countries[i]
                p_title.font.name = 'Knewave'
                p_title.font.size = Pt(40)
                p_title.font.color.rgb = RGBColor(218,57,21)

                

                # Set the zorder of the title shape to ensure it's above the background image
                txBox_title.zorder = 10

                # Add attractions with pointers
                left = Inches(0.8)
                top = Inches(1.25)
                width = Inches(4.5)
                height = Inches(5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True

                # Calculate the length of the text box for the body text
                body_text_length = Inches(11.3)

                # Resize the body text font size if it exceeds the width
                for j, attraction in enumerate(nested_list[i]):
                    p = tf.add_paragraph()
                    p.text = f"{j+1}. {attraction}"
                    p.font.size = Pt(16) if txBox.width > body_text_length else Pt(20)  # Adjust font size
                    p.space_before = Pt(14)  # Adjust spacing between each attraction
                    p.font.color.rgb = RGBColor(139,50,28)

                # Add image from URL
                image = image_urls[i]
                try:
                    response = requests.get(image)
                    if response.status_code == 200:
                        image_stream = BytesIO(response.content)
                        slide.shapes.add_picture(image_stream, Inches(6), Inches(2.5), width=Inches(3), height=Inches(3))
                except Exception as e:
                    print(f"Failed to load image from {image}: {e}") 
        end_slide_layout = prs.slide_layouts[0]  # Layout index 0 is the title slide layout
        end_slide = prs.slides.add_slide(end_slide_layout)


        # Add background image
        try:
            bg_image_path = "themes/fall/fall-conclusion.png"  # Path to your local image
            bg_image = end_slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            bg_image.zorder = -10  # Set the zorder to place it behind other shapes
        except Exception as e:
            print(f"Failed to load background image: {e}")

        left = Inches(2)  # Adjust position as needed
        top = Inches(2)  # Adjust position as needed
        width = Inches(8)  # Adjust size as needed
        height = Inches(1)  # Adjust size as needed
        txBox_end = end_slide.shapes.add_textbox(left, top, width, height)
        tf_end = txBox_end.text_frame
        p_end = tf_end.add_paragraph()
        p_end.text = 'Thank You!'
        p_end.font.bold = True
        p_end.font.size = Pt(55) 
        p_end.font.color.rgb = RGBColor(139,50,28)



    ## slides formatting for spring
    if theme == 'spring':
        title_slide_layout = prs.slide_layouts[0]  # Layout index 0 is the title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)

        # Add background image
        try:
            bg_image_path = "themes/spring/spring-intro.png"  # Path to your local image
            bg_image = title_slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            bg_image.zorder = -10  # Set the zorder to place it behind other shapes
        except Exception as e:
            print(f"Failed to load background image: {e}")

        left = Inches(1)  # Adjust position as needed
        top = Inches(1)  # Adjust position as needed
        width = Inches(8)  # Adjust size as needed
        height = Inches(3)  # Adjust size as needed
        txBox_intro = title_slide.shapes.add_textbox(left, top, width, height)
        tf_intro = txBox_intro.text_frame
        p_intro = tf_intro.add_paragraph()
        p_intro.text = title
        tf_intro.word_wrap = True
        p_intro.font.bold = True
        p_intro.font.size = Pt(48)
        p_intro.font.color.rgb = RGBColor(63,102,47)
        for paragraph in tf_intro.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER


        for i in range(num_slides):
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use the layout for title slides (layout index 5)

                try:
                    bg_image_path = "themes/spring/spring-bg.png"  # Path to your local image
                    bg_image = slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
                    bg_image.zorder = -10  # Set the zorder to place it behind other shapes

                except Exception as e:
                    print(f"Failed to load background image: {e}")



                # Add title as a text box with similar structure to attractions with pointers
                left = Inches(4)
                top = Inches(0.1)
                width = Inches(9)
                height = Inches(1)
                txBox_title = slide.shapes.add_textbox(left, top, width, height)
                tf_title = txBox_title.text_frame
                tf_title.word_wrap = True
                p_title = tf_title.add_paragraph()
                p_title.text = countries[i]
                p_title.font.name = 'Knewave'
                p_title.font.size = Pt(40)
                p_title.font.color.rgb = RGBColor(63,102,47)

                

                # Set the zorder of the title shape to ensure it's above the background image
                txBox_title.zorder = 10

                # Add attractions with pointers
                left = Inches(0.8)
                top = Inches(1.25)
                width = Inches(4.5)
                height = Inches(5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True

                # Calculate the length of the text box for the body text
                body_text_length = Inches(11.3)

                # Resize the body text font size if it exceeds the width
                for j, attraction in enumerate(nested_list[i]):
                    p = tf.add_paragraph()
                    p.text = f"{j+1}. {attraction}"
                    p.font.size = Pt(16) if txBox.width > body_text_length else Pt(20)  # Adjust font size
                    p.space_before = Pt(14)  # Adjust spacing between each attraction

                # Add image from URL
                image = image_urls[i]
                try:
                    response = requests.get(image)
                    if response.status_code == 200:
                        image_stream = BytesIO(response.content)
                        slide.shapes.add_picture(image_stream, Inches(6), Inches(2.5), width=Inches(3), height=Inches(3))
                except Exception as e:
                    print(f"Failed to load image from {image}: {e}")

        end_slide_layout = prs.slide_layouts[0]  # Layout index 0 is the title slide layout
        end_slide = prs.slides.add_slide(end_slide_layout)


        # Add background image
        try:
            bg_image_path = "themes/spring/spring-conclusion.png"  # Path to your local image
            bg_image = end_slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            bg_image.zorder = -10  # Set the zorder to place it behind other shapes
        except Exception as e:
            print(f"Failed to load background image: {e}")

        left = Inches(3)  # Adjust position as needed
        top = Inches(3)  # Adjust position as needed
        width = Inches(8)  # Adjust size as needed
        height = Inches(1)  # Adjust size as needed
        txBox_end = end_slide.shapes.add_textbox(left, top, width, height)
        tf_end = txBox_end.text_frame
        p_end = tf_end.add_paragraph()
        p_end.text = 'Thank You!'
        p_end.font.bold = True
        p_end.font.size = Pt(60)
        p_end.font.color.rgb = RGBColor(63,102,47)


    # Save the presentation
    prs.save('themed_presentation.pptx')
    return prs